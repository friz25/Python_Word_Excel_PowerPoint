"""3) Чтение PowerPoint файлов"""
import collections
import collections.abc
from pptx import Presentation

root = Presentation('example2.pptx')

file = open('text.txt', 'w')

for slide in root.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                print(run.text)
                file.write(run.text + '\n')

file.close()