"""2) Добавление абзацев и шрифтов"""
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt

root = Presentation()

blank_slide = root.slide_layouts[6]

slide = root.slides.add_slide(blank_slide)

left = top = width = height = Inches(1)

txBox = slide.shapes.add_textbox(left, top, width, height)

tf = txBox.text_frame
"""1й текст"""
tf.text = 'Текст 1'
"""2й текст"""
p = tf.add_paragraph()

p.text= 'Текст 2'

p.font.bold = True
p.font.italic = True
"""3й текст"""
p = tf.add_paragraph()
p.text = 'Текст 3'

p.font.size = Pt(40)
p.font.name = 'Arial'

root.save('example2.pptx')