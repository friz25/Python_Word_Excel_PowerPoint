"""1) Создание слайдов"""
import collections
import collections.abc
from pptx import Presentation

root = Presentation()
"""Первый слайд"""
first_slide_layout = root.slide_layouts[4]

slide = root.slides.add_slide(first_slide_layout)

slide.shapes.title.text = 'Тестовый слайд'

slide.placeholders[1].text = 'Подзаголовок тестового слайда'
slide.placeholders[3].text = 'Подзаголовок тестового слайда 2'
"""Второй слайд"""
first_slide_layout2 = root.slide_layouts[2]

slide = root.slides.add_slide(first_slide_layout2)

slide.shapes.title.text = 'Тестовый слайд'

slide.placeholders[1].text = 'Подзаголовок тестового слайда'


root.save('example1.pptx')